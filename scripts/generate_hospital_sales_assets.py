from pathlib import Path

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PPTInches, Pt as PPTPt


ROOT = Path(__file__).resolve().parents[1]
DOCX_DIR = ROOT / "downloads" / "docx"
PPTX_DIR = ROOT / "downloads" / "pptx"

BRAND_BLUE = RGBColor(22, 78, 99)
BRAND_TEAL = RGBColor(8, 145, 178)
BRAND_ORANGE = RGBColor(234, 88, 12)
TEXT_DARK = RGBColor(30, 41, 59)
TEXT_MUTED = RGBColor(71, 85, 105)

PPT_BLUE = PPTColor(22, 78, 99)
PPT_TEAL = PPTColor(8, 145, 178)
PPT_ORANGE = PPTColor(234, 88, 12)
PPT_BG = PPTColor(248, 250, 252)
PPT_DARK = PPTColor(15, 23, 42)
PPT_MUTED = PPTColor(71, 85, 105)


CASES = [
    {
        "slug": "企业微信智慧就医案例",
        "title": "某市人民医院门诊服务案例",
        "subtitle": "企业微信智慧就医服务实践",
        "cover_subtitle": "提升门诊服务效率与患者满意度",
        "filename": "微语AI医院案例-企业微信智慧就医",
        "cover_quote": "把门诊咨询前移到企业微信，让患者在到院前就完成导诊、挂号和就诊准备。",
        "summary": [
            "某市人民医院设有 3 个门诊院区，工作日日均门诊量超过 1.2 万人次，导诊与挂号咨询长期高频集中。",
            "医院以企业微信作为患者服务统一入口，建设 7x24 小时智慧就医服务体系，覆盖导诊、挂号、检查准备、候诊提醒和人工协同。",
            "项目目标是把大量重复咨询前移到线上处理，改善门诊现场秩序，并形成可复制的门诊服务标准流程。",
        ],
        "pain_points": [
            "门诊大厅、热线和公众号同时承接患者咨询，高峰时段重复问答密集，导医台与热线压力明显。",
            "不同院区科室分布、挂号要求和检查准备事项复杂，患者在到院后仍需反复确认流程。",
            "医院已有企业微信协同基础，但患者服务入口分散，缺少统一记录和连续服务机制。",
            "院方希望先从门诊服务中心与重点专科切入，验证智慧就医服务对门诊秩序和患者体验的改善效果。",
        ],
        "scenes": [
            [
                "门诊导诊分流",
                "患者在企业微信咨询“头晕挂什么科”“儿童发热去哪里看”时，系统根据主诉、年龄和就诊需求推荐院区、科室和就诊建议。",
            ],
            [
                "挂号与检查前提醒",
                "患者确认就诊后，系统自动发送挂号方式、证件要求、空腹检查提醒、医保说明等内容，减少窗口反复解释。",
            ],
            [
                "候诊与院内流程提醒",
                "围绕取号、候诊、缴费、检验检查、报告打印等节点，企业微信主动推送流程提示，降低患者焦虑和现场问询量。",
            ],
            [
                "复杂问题转人工",
                "对于多病史、跨科室、特殊检查安排等复杂问题，系统自动转给导医或医助接管，保证服务连续性。",
            ],
        ],
        "solution_modules": [
            "企业微信统一服务入口：将导诊咨询、挂号引导、院内流程说明统一到一个患者服务窗口。",
            "医院知识库：沉淀科室介绍、挂号规则、检查须知、院区地图、医保与就诊说明。",
            "智能工作流：识别高频咨询后自动推送标准答案、流程卡片或人工转接。",
            "服务数据看板：持续统计咨询量、自动响应率、人工接管率和重点问题分布。",
        ],
        "sales_points": [
            "门诊高频咨询由系统自动承接后，导医台与热线的重复答复压力明显下降。",
            "患者在到院前即可完成挂号、检查准备和院区路径了解，现场排队和问询拥堵得到缓解。",
            "复杂问题由系统快速转人工处理，保证了服务效率，也保留了人工专业判断。",
            "项目先在门诊服务中心和两个重点专科试点，成熟后逐步复制到多院区与更多门诊场景。",
        ],
        "roi": [
            ["导诊咨询自动化率", "65%-80%"],
            ["人工客服压力下降", "30%-50%"],
            ["预约挂号转化提升", "20%-35%"],
            ["患者满意度提升", "15%-25%"],
            ["试点上线周期", "2-4 周"],
        ],
        "buyer_map": [
            ["门诊服务部", "负责梳理导诊问答口径、优化患者接待流程，并参与服务规则设计。"],
            ["互联网医院/运营中心", "负责企业微信入口建设、消息触达和线上服务闭环联动。"],
            ["信息科", "负责账号接入、数据安全、接口联调和系统运行保障。"],
            ["重点专科/医助团队", "负责复杂问题人工接管和重点场景试点验证。"],
        ],
        "landing_steps": [
            "第 1 周：梳理门诊高频咨询、整理科室与挂号规则、确定试点科室和服务团队。",
            "第 2 周：完成企业微信接入、知识库配置、导诊流程和人工转接规则上线。",
            "第 3 周：联调挂号与流程提醒能力，开展导医与客服培训并灰度试运行。",
            "第 4 周：复盘咨询数据和患者反馈，补充知识库并规划下一阶段复制范围。",
        ],
        "demo_script": [
            "患者在企业微信咨询“胃痛挂哪个科”，系统根据主诉推荐消化内科并说明院区位置。",
            "患者继续询问“今天还能挂号吗”，系统返回挂号方式、剩余时段和就诊提醒。",
            "首次就诊患者收到检查准备、缴费方式、候诊节点等流程卡片，减少现场重复问询。",
            "当患者进一步提出复杂病史问题时，会话自动转给导医或医助继续服务。",
        ],
    },
    {
        "slug": "企业微信健康咨询案例",
        "title": "某慢病专科医院随访服务案例",
        "subtitle": "企业微信健康咨询与随访实践",
        "cover_subtitle": "提升复诊执行率与持续服务体验",
        "filename": "微语AI医院案例-企业微信健康咨询",
        "cover_quote": "把离院后的患者咨询、提醒与回访统一到企业微信，形成连续服务闭环。",
        "summary": [
            "某慢病专科医院重点服务糖尿病、高血压和术后康复患者，复诊患者占比高，离院后咨询需求持续存在。",
            "医院以企业微信为患者触达主渠道，围绕健康咨询、报告说明、用药提醒和复诊随访构建连续服务闭环。",
            "项目重点是把一次门诊服务延伸为长期健康管理，提升复诊执行率和重点患者跟进效率。",
        ],
        "pain_points": [
            "患者离院后会持续咨询饮食、用药、复查时间和康复注意事项，人工随访工作量持续增长。",
            "慢病复诊提醒、体检异常解释和术后康复建议分散在电话、微信群和线下回访中，难以统一管理。",
            "医护和客服团队重复答复相似问题，服务记录不连续，重点患者分层管理推进缓慢。",
            "医院希望在不显著增加人力的前提下，建立规范化、可持续的健康咨询与随访服务体系。",
        ],
        "scenes": [
            [
                "离院后健康咨询",
                "患者在企业微信咨询“术后多久能正常饮食”“血糖偏高要注意什么”时，系统提供标准化健康建议和注意事项。",
            ],
            [
                "报告解读与风险提醒",
                "患者发送体检或检验结果后，系统进行非诊断式说明，并根据风险等级提醒人工跟进或复查。",
            ],
            [
                "用药与复诊提醒",
                "系统根据病种与治疗阶段，在企业微信中自动推送用药提醒、复查时间、复诊建议和随访问卷。",
            ],
            [
                "重点患者分层管理",
                "医院按病种、阶段、风险和服务状态为患者打标签，便于重点患者持续跟进。",
            ],
        ],
        "solution_modules": [
            "企业微信患者服务台：沉淀患者档案、会话记录、标签和服务轨迹。",
            "健康咨询知识库：统一病种说明、术后康复、饮食建议、用药提醒和体检解释口径。",
            "随访工作流引擎：自动安排提醒、问卷、人工回访和异常预警流程。",
            "运营数据看板：查看咨询触达率、复诊转化率、重点患者跟进情况和活跃度。",
        ],
        "sales_points": [
            "常见健康咨询和随访提醒由系统自动承接后，护士与客服团队的重复性工作明显减少。",
            "患者在离院后仍能持续获得标准化服务和复诊提醒，复查与复诊执行率更稳定。",
            "重点患者可被及时识别并转入人工随访，提高专病管理和康复管理连续性。",
            "项目上线后，医院逐步建立了可量化、可持续优化的患者健康管理服务机制。",
        ],
        "roi": [
            ["常见咨询自动响应率", "70%-85%"],
            ["护士/客服随访工作量下降", "35%-55%"],
            ["复诊邀约转化提升", "15%-30%"],
            ["患者持续活跃提升", "20%-40%"],
            ["慢病管理项目回本周期", "3-6 个月"],
        ],
        "buyer_map": [
            ["互联网医院/运营中心", "负责企业微信患者触达、线上服务运营和患者留存管理。"],
            ["护理部/随访中心", "负责随访模板梳理、异常患者跟进和服务规范执行。"],
            ["体检中心/康复中心", "负责重点患者场景导入、内容维护和长期服务设计。"],
            ["信息科与合规部门", "负责账号接入、数据安全、日志留痕和系统稳定运行。"],
        ],
        "landing_steps": [
            "第 1 周：明确试点病种与患者分组，导入健康问答和随访模板。",
            "第 2 周：完成企业微信账号、标签体系、提醒策略与回访流程配置。",
            "第 3 周：启动患者分层触达与人工协同闭环，观察触达率和回复效果。",
            "第 4 周：根据患者反馈与复诊数据优化内容，规划下一阶段扩展场景。",
        ],
        "demo_script": [
            "患者在企业微信询问“胆囊术后饮食怎么安排”，系统立即给出标准化康复建议。",
            "患者发送体检异常指标后，系统进行非诊断式说明，并提示是否需要进一步复查。",
            "在复诊时间点到来前，系统自动提醒患者复诊并引导完成预约。",
            "对高风险或高关注患者，系统自动提醒人工团队继续回访与跟进。",
        ],
    },
]


def ensure_output_dirs() -> None:
    DOCX_DIR.mkdir(parents=True, exist_ok=True)
    PPTX_DIR.mkdir(parents=True, exist_ok=True)


def set_cell_text(cell, text, bold=False):
    cell.text = ""
    paragraph = cell.paragraphs[0]
    run = paragraph.add_run(text)
    run.bold = bold
    run.font.name = "PingFang SC"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "PingFang SC")
    run.font.size = Pt(10.5)


def shade_cell(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def add_bullets(doc: Document, items, level=0):
    for item in items:
        paragraph = doc.add_paragraph(style="List Bullet")
        if level:
            paragraph.paragraph_format.left_indent = Inches(0.3 * level)
        run = paragraph.add_run(item)
        run.font.name = "PingFang SC"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "PingFang SC")
        run.font.size = Pt(11)
        run.font.color.rgb = TEXT_DARK


def add_title(doc: Document, text: str, size: int, color: RGBColor):
    paragraph = doc.add_paragraph()
    paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = paragraph.add_run(text)
    run.bold = True
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = "PingFang SC"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "PingFang SC")


def add_heading(doc: Document, text: str, level: int = 1):
    paragraph = doc.add_paragraph()
    run = paragraph.add_run(text)
    run.bold = True
    run.font.name = "PingFang SC"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "PingFang SC")
    run.font.size = Pt(16 if level == 1 else 13)
    run.font.color.rgb = BRAND_BLUE if level == 1 else BRAND_TEAL


def add_intro_block(doc: Document, quote: str, summary):
    quote_para = doc.add_paragraph()
    quote_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    quote_run = quote_para.add_run(quote)
    quote_run.italic = True
    quote_run.font.size = Pt(13)
    quote_run.font.name = "PingFang SC"
    quote_run._element.rPr.rFonts.set(qn("w:eastAsia"), "PingFang SC")
    quote_run.font.color.rgb = TEXT_MUTED
    add_bullets(doc, summary)


def add_two_col_table(doc: Document, title: str, headers, rows, fill="E0F2FE"):
    add_heading(doc, title)
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    header_cells = table.rows[0].cells
    for index, header in enumerate(headers):
        set_cell_text(header_cells[index], header, bold=True)
        shade_cell(header_cells[index], fill)
    for row in rows:
        cells = table.add_row().cells
        for index, value in enumerate(row):
            set_cell_text(cells[index], value)
    doc.add_paragraph()


def create_docx(case: dict) -> Path:
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.8)
    section.bottom_margin = Inches(0.7)
    section.left_margin = Inches(0.9)
    section.right_margin = Inches(0.9)

    styles = doc.styles
    styles["Normal"].font.name = "PingFang SC"
    styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "PingFang SC")
    styles["Normal"].font.size = Pt(11)

    add_title(doc, case["title"], 22, BRAND_BLUE)
    add_title(doc, case["subtitle"], 17, BRAND_TEAL)
    add_title(doc, case["cover_subtitle"], 12, BRAND_ORANGE)
    doc.add_paragraph()
    add_intro_block(doc, case["cover_quote"], case["summary"])

    add_heading(doc, "一、项目背景")
    add_bullets(doc, case["pain_points"])

    add_two_col_table(doc, "二、医院实际应用场景", ["场景", "实际应用说明"], case["scenes"])

    add_heading(doc, "三、建设内容")
    add_bullets(doc, case["solution_modules"])

    add_heading(doc, "四、上线成效")
    add_bullets(doc, case["sales_points"])

    add_two_col_table(doc, "五、试点阶段成效指标", ["指标", "预期结果"], case["roi"], fill="DCFCE7")
    add_two_col_table(doc, "六、院内协同方式", ["部门", "说明"], case["buyer_map"], fill="FFEDD5")

    add_heading(doc, "七、30 天落地步骤")
    add_bullets(doc, case["landing_steps"])

    add_heading(doc, "八、典型患者服务流程")
    add_bullets(doc, case["demo_script"])

    output = DOCX_DIR / f"{case['filename']}.docx"
    doc.save(output)
    return output


def add_slide_title(slide, title, subtitle=None, accent=PPT_BLUE):
    title_box = slide.shapes.add_textbox(PPTInches(0.55), PPTInches(0.35), PPTInches(8.2), PPTInches(0.7))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = PPTPt(24)
    title_para.font.bold = True
    title_para.font.color.rgb = accent
    title_para.font.name = "Arial"
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(PPTInches(0.58), PPTInches(0.95), PPTInches(8.4), PPTInches(0.45))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = PPTPt(10.5)
        subtitle_para.font.color.rgb = PPT_MUTED
        subtitle_para.font.name = "Arial"


def add_banner(slide, color=PPT_TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, PPTInches(0), PPTInches(0), PPTInches(13.33), PPTInches(0.18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_bullet_box(slide, x, y, w, h, items, color=PPT_DARK, font_size=18):
    box = slide.shapes.add_textbox(x, y, w, h)
    frame = box.text_frame
    frame.word_wrap = True
    first = True
    for item in items:
        para = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        para.text = item
        para.level = 0
        para.font.size = PPTPt(font_size)
        para.font.color.rgb = color
        para.font.name = "Arial"
        para.space_after = PPTPt(8)
        para.bullet = True


def add_metric_cards(slide, metrics):
    positions = [0.7, 3.45, 6.2, 8.95, 11.0]
    widths = [2.4, 2.4, 2.4, 1.8, 1.8]
    for index, (name, value) in enumerate(metrics):
        left = PPTInches(positions[index])
        top = PPTInches(1.65)
        width = PPTInches(widths[index])
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, PPTInches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = PPT_BG
        card.line.color.rgb = PPT_TEAL if index % 2 == 0 else PPT_ORANGE

        title = slide.shapes.add_textbox(left + PPTInches(0.15), top + PPTInches(0.18), width - PPTInches(0.3), PPTInches(0.35))
        tf = title.text_frame.paragraphs[0]
        tf.text = name
        tf.font.size = PPTPt(11)
        tf.font.color.rgb = PPT_MUTED
        tf.font.name = "Arial"

        value_box = slide.shapes.add_textbox(left + PPTInches(0.15), top + PPTInches(0.6), width - PPTInches(0.3), PPTInches(0.55))
        vf = value_box.text_frame.paragraphs[0]
        vf.text = value
        vf.font.size = PPTPt(21)
        vf.font.bold = True
        vf.font.color.rgb = PPT_DARK
        vf.font.name = "Arial"


def add_two_column_table_like(slide, title, rows, left_title="模块", right_title="说明"):
    add_slide_title(slide, title)
    left = PPTInches(0.7)
    top = PPTInches(1.55)
    width = PPTInches(11.9)
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, PPTInches(0.45))
    header.fill.solid()
    header.fill.fore_color.rgb = PPT_BLUE
    header.line.fill.background()
    for offset, text in [(0.25, left_title), (5.0, right_title)]:
        box = slide.shapes.add_textbox(left + PPTInches(offset), top + PPTInches(0.08), PPTInches(3.5), PPTInches(0.22))
        para = box.text_frame.paragraphs[0]
        para.text = text
        para.font.size = PPTPt(12)
        para.font.bold = True
        para.font.color.rgb = PPTColor(255, 255, 255)
        para.font.name = "Arial"
    row_top = top + PPTInches(0.5)
    for idx, row in enumerate(rows):
        bg = PPTColor(255, 255, 255) if idx % 2 == 0 else PPTColor(241, 245, 249)
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, row_top, width, PPTInches(0.72))
        band.fill.solid()
        band.fill.fore_color.rgb = bg
        band.line.color.rgb = PPTColor(226, 232, 240)
        for offset, text, font_color in [(0.22, row[0], PPT_BLUE), (5.0, row[1], PPT_DARK)]:
            box = slide.shapes.add_textbox(left + PPTInches(offset), row_top + PPTInches(0.08), PPTInches(5.2), PPTInches(0.55))
            para = box.text_frame.paragraphs[0]
            para.text = text
            para.font.size = PPTPt(14 if offset < 1 else 12.5)
            para.font.bold = offset < 1
            para.font.color.rgb = font_color
            para.font.name = "Arial"
        row_top += PPTInches(0.75)


def add_cover_slide(prs: Presentation, case: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = PPTColor(240, 249, 255)
    background.line.fill.background()
    add_banner(slide, PPT_TEAL)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(0.65), PPTInches(0.85), PPTInches(12.0), PPTInches(5.7))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PPTColor(255, 255, 255)
    accent.line.color.rgb = PPTColor(186, 230, 253)

    title = slide.shapes.add_textbox(PPTInches(0.95), PPTInches(1.2), PPTInches(10.8), PPTInches(1.5))
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = case["subtitle"]
    p.font.size = PPTPt(28)
    p.font.bold = True
    p.font.color.rgb = PPT_BLUE
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = case["cover_subtitle"]
    p2.font.size = PPTPt(16)
    p2.font.color.rgb = PPT_ORANGE
    p2.font.name = "Arial"

    quote = slide.shapes.add_textbox(PPTInches(0.98), PPTInches(2.35), PPTInches(10.4), PPTInches(0.8))
    q = quote.text_frame.paragraphs[0]
    q.text = case["cover_quote"]
    q.font.size = PPTPt(18)
    q.font.color.rgb = PPT_MUTED
    q.font.name = "Arial"

    add_bullet_box(slide, PPTInches(1.0), PPTInches(3.15), PPTInches(10.2), PPTInches(2.5), case["summary"], font_size=18)


def add_pain_slide(prs: Presentation, case: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide)
    add_slide_title(slide, "项目背景", "结合医院现有服务现状与建设目标")
    left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(0.7), PPTInches(1.55), PPTInches(5.85), PPTInches(5.2))
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = PPTColor(255, 255, 255)
    left_panel.line.color.rgb = PPTColor(186, 230, 253)
    right_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(6.8), PPTInches(1.55), PPTInches(5.85), PPTInches(5.2))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = PPTColor(255, 247, 237)
    right_panel.line.color.rgb = PPTColor(254, 215, 170)
    add_bullet_box(slide, PPTInches(0.95), PPTInches(1.9), PPTInches(5.2), PPTInches(4.4), case["pain_points"][:2], font_size=17)
    add_bullet_box(slide, PPTInches(7.05), PPTInches(1.9), PPTInches(5.15), PPTInches(4.4), case["pain_points"][2:], font_size=17)


def add_scene_slide(prs: Presentation, case: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, PPT_ORANGE)
    add_two_column_table_like(slide, "医院实际应用场景", case["scenes"], "场景", "实际应用说明")


def add_modules_slide(prs: Presentation, case: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, PPT_TEAL)
    add_slide_title(slide, "建设内容", "围绕企业微信患者服务入口构建能力")
    add_bullet_box(slide, PPTInches(0.78), PPTInches(1.65), PPTInches(5.9), PPTInches(4.9), case["solution_modules"], font_size=16)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(7.1), PPTInches(1.75), PPTInches(5.35), PPTInches(4.65))
    box.fill.solid()
    box.fill.fore_color.rgb = PPTColor(239, 246, 255)
    box.line.color.rgb = PPTColor(125, 211, 252)
    add_bullet_box(slide, PPTInches(7.35), PPTInches(2.05), PPTInches(4.8), PPTInches(4.1), case["sales_points"], font_size=16)


def add_roi_slide(prs: Presentation, case: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, PPT_BLUE)
    add_slide_title(slide, "试点阶段成效", "通过数据验证服务效率与患者体验改善效果")
    add_metric_cards(slide, case["roi"])
    note = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(0.8), PPTInches(3.7), PPTInches(12.0), PPTInches(2.35))
    note.fill.solid()
    note.fill.fore_color.rgb = PPTColor(255, 255, 255)
    note.line.color.rgb = PPTColor(186, 230, 253)
    add_bullet_box(slide, PPTInches(1.05), PPTInches(4.0), PPTInches(11.5), PPTInches(1.8), [
        "建议试点阶段重点关注自动化响应率、人工工作量下降和患者满意度提升。",
        "项目可按单院区、单专科或单中心先行落地，逐步扩展到更多服务场景。",
        "通过阶段性复盘持续完善知识库、流程规则和人工协同机制。",
    ], font_size=15)


def add_buyer_slide(prs: Presentation, case: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, PPT_TEAL)
    add_two_column_table_like(slide, "院内协同方式", case["buyer_map"], "部门", "协同方式")


def add_closing_slide(prs: Presentation, case: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, PPT_ORANGE)
    add_slide_title(slide, "案例亮点", "聚焦医院实际应用效果与患者服务流程")
    left_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(0.75), PPTInches(1.55), PPTInches(4.0), PPTInches(5.25))
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = PPTColor(255, 255, 255)
    left_shape.line.color.rgb = PPTColor(186, 230, 253)
    mid_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(4.95), PPTInches(1.55), PPTInches(3.7), PPTInches(5.25))
    mid_shape.fill.solid()
    mid_shape.fill.fore_color.rgb = PPTColor(240, 253, 244)
    mid_shape.line.color.rgb = PPTColor(134, 239, 172)
    right_shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, PPTInches(8.85), PPTInches(1.55), PPTInches(3.75), PPTInches(5.25))
    right_shape.fill.solid()
    right_shape.fill.fore_color.rgb = PPTColor(255, 247, 237)
    right_shape.line.color.rgb = PPTColor(254, 215, 170)

    for x, title in [(1.0, "上线成效"), (5.2, "患者流程"), (9.1, "推进特点")]:
        box = slide.shapes.add_textbox(PPTInches(x), PPTInches(1.8), PPTInches(2.8), PPTInches(0.3))
        para = box.text_frame.paragraphs[0]
        para.text = title
        para.font.size = PPTPt(18)
        para.font.bold = True
        para.font.color.rgb = PPT_BLUE
        para.font.name = "Arial"

    add_bullet_box(slide, PPTInches(1.0), PPTInches(2.2), PPTInches(3.4), PPTInches(4.2), case["sales_points"][:3], font_size=14)
    add_bullet_box(slide, PPTInches(5.2), PPTInches(2.2), PPTInches(3.1), PPTInches(4.2), case["demo_script"][:3], font_size=14)
    add_bullet_box(slide, PPTInches(9.1), PPTInches(2.2), PPTInches(3.0), PPTInches(4.2), case["landing_steps"][:3], font_size=14)


def add_landing_slide(prs: Presentation, case: dict):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_banner(slide, PPT_BLUE)
    add_slide_title(slide, "30 天试点落地路径", "帮助医院快速完成试点建设与阶段复盘")
    add_bullet_box(slide, PPTInches(0.85), PPTInches(1.7), PPTInches(12.0), PPTInches(4.8), case["landing_steps"], font_size=20)
    footer = slide.shapes.add_textbox(PPTInches(0.9), PPTInches(6.4), PPTInches(11.6), PPTInches(0.3))
    para = footer.text_frame.paragraphs[0]
    para.text = "建议结合试点验收指标与阶段复盘机制，稳步推进后续复制与扩展。"
    para.font.size = PPTPt(12)
    para.font.color.rgb = PPT_MUTED
    para.font.name = "Arial"
    para.alignment = PP_ALIGN.CENTER


def create_pptx(case: dict) -> Path:
    prs = Presentation()
    prs.slide_width = PPTInches(13.333)
    prs.slide_height = PPTInches(7.5)
    add_cover_slide(prs, case)
    add_pain_slide(prs, case)
    add_scene_slide(prs, case)
    add_modules_slide(prs, case)
    add_roi_slide(prs, case)
    add_buyer_slide(prs, case)
    add_closing_slide(prs, case)
    add_landing_slide(prs, case)
    output = PPTX_DIR / f"{case['filename']}.pptx"
    prs.save(output)
    return output


def main() -> None:
    ensure_output_dirs()
    for case in CASES:
        docx_path = create_docx(case)
        pptx_path = create_pptx(case)
        print(f"generated:{docx_path}")
        print(f"generated:{pptx_path}")


if __name__ == "__main__":
    main()